import collections
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    # Use blank slide layout for all to have full control
    blank_layout = prs.slide_layouts[6] 

    # ----- Slide 1: Title Page -----
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Title
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "HONEYWELL HACKATHON 2026\nTITLE PAGE"
    p.font.bold = True
    p.font.size = Pt(28)
    
    # Content
    contentBox = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
    tf2 = contentBox.text_frame
    tf2.word_wrap = True
    points = [
        "• Problem Statement ID: (Enter PS ID)",
        "• Problem Statement Title: AI-Powered Behavioral Anomaly Detection",
        "• Theme: Cybersecurity & Smart Automation",
        "• PS Category: Software",
        "• Student Name: Atharv Sinha",
        "• Student ID: (Enter Student ID)"
    ]
    for point in points:
        p = tf2.add_paragraph()
        p.text = point
        p.font.size = Pt(20)

    # ----- Slide 2: Proposed Solution -----
    slide2 = prs.slides.add_slide(blank_layout)
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "HONEYWELL HACKATHON 2026\nAI-Driven Behavioral Anomaly Detection"
    p.font.bold = True
    p.font.size = Pt(24)

    contentBox = slide2.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    tf2 = contentBox.text_frame
    tf2.word_wrap = True
    points = [
        "Proposed Solution:",
        "• Detailed explanation of the proposed solution:",
        "  - AI/ML system that models normal access behavior for users and IoT devices.",
        "  - Automatically detects and classifies 8 intrusion types (e.g., brute force, lateral movement).",
        "  - Provides explainable risk scores through an analyst-facing SOC dashboard.",
        "",
        "• How it addresses the problem:",
        "  - Replaces rigid static rules with dynamic behavioral baselines.",
        "  - Mitigates alert fatigue by intelligently scoring and prioritizing threats.",
        "",
        "• Innovation and uniqueness of the solution:",
        "  - Dual Pipeline: Isolation Forest (unsupervised profiling) + Random Forest (supervised multi-class).",
        "  - Explainability Engine: Translates ML decisions into human-readable alerts.",
        "  - Cold-start handling for new entities using statistical profiling fallbacks."
    ]
    for i, point in enumerate(points):
        p = tf2.add_paragraph()
        p.text = point
        if i == 0 or point.startswith("•"):
            p.font.bold = True
        p.font.size = Pt(16)


    # ----- Slide 3: Technical Approach -----
    slide3 = prs.slides.add_slide(blank_layout)
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "HONEYWELL HACKATHON 2026\nTECHNICAL APPROACH"
    p.font.bold = True
    p.font.size = Pt(24)

    contentBox = slide3.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    tf2 = contentBox.text_frame
    tf2.word_wrap = True
    points = [
        "• Technologies to be used:",
        "  - Machine Learning: Python, scikit-learn, Pandas, NumPy.",
        "  - Visualization & UI: Plotly Dash, Flask, Bootstrap.",
        "  - Hardware: Standard Laptop/Server (Designed for lightweight, local execution).",
        "",
        "• Methodology and process for implementation:",
        "  1. Data Generator Layer: Synthetic log generation simulating 50K+ events and 8 attack patterns.",
        "  2. ML Pipeline - Profiling: Isolation Forest isolates 'normal' entity behavior.",
        "  3. ML Pipeline - Classification: Random Forest determines specific attack types.",
        "  4. Explainability Layer: Feature attribution maps statistical anomalies to human-readable text.",
        "  5. Dashboard UI: Real-time Plotly Dash interface visualizes KPIs, geo-maps, and alert queues."
    ]
    for i, point in enumerate(points):
        p = tf2.add_paragraph()
        p.text = point
        if point.startswith("•"):
            p.font.bold = True
        p.font.size = Pt(16)


    # ----- Slide 4: Feasibility and Viability -----
    slide4 = prs.slides.add_slide(blank_layout)
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "HONEYWELL HACKATHON 2026\nFEASIBILITY AND VIABILITY"
    p.font.bold = True
    p.font.size = Pt(24)

    contentBox = slide4.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    tf2 = contentBox.text_frame
    tf2.word_wrap = True
    points = [
        "• Analysis of the feasibility of the idea:",
        "  - Cost-Effective: Built entirely on an open-source Python stack (No vendor lock-in).",
        "  - Highly Performant: The 50K event pipeline trains and infers in under 30 seconds.",
        "  - Extensible: Modular architecture easily integrates with standard SIEM APIs (Splunk, Elastic).",
        "",
        "• Potential challenges and risks:",
        "  - Evolving Threats (Concept Drift): Attack patterns change over time.",
        "  - Alert Fatigue: Improperly tuned models can overwhelm SOC analysts with false positives.",
        "  - Big Data Scaling: Processing millions of logs daily requires distributed computing.",
        "",
        "• Strategies for overcoming these challenges:",
        "  - Drift Mitigation: Implement continuous, scheduled model retraining.",
        "  - Analyst Feedback Loop: Tunable contamination thresholds per-entity to reduce noise.",
        "  - Scaling Roadmap: Transition MVP architecture to Apache Spark / Kafka for enterprise deployment."
    ]
    for i, point in enumerate(points):
        p = tf2.add_paragraph()
        p.text = point
        if point.startswith("•"):
            p.font.bold = True
        p.font.size = Pt(16)


    # ----- Slide 5: Artifacts -----
    slide5 = prs.slides.add_slide(blank_layout)
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "HONEYWELL HACKATHON 2026\nARTIFACTS"
    p.font.bold = True
    p.font.size = Pt(24)

    contentBox = slide5.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1.5))
    tf2 = contentBox.text_frame
    tf2.word_wrap = True
    points = [
        "• Relevant artifacts:",
        "  - GitHub Repository: github.com/AtharvSinha26052005/honeywell_AI_Anomoly_Detection",
        "  - Prototype Output: 48,700 Events, 97.45% Accuracy across 8 attack patterns.",
        "  - Snaps of the Solution: Real-time dashboard components below."
    ]
    for point in points:
        p = tf2.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        if point.startswith("•"): p.font.bold = True

    # Add images
    try:
        slide5.shapes.add_picture("dashboard_top.png", Inches(0.5), Inches(2.5), height=Inches(2.2))
        slide5.shapes.add_picture("dashboard_middle.png", Inches(5.0), Inches(2.5), height=Inches(2.2))
        slide5.shapes.add_picture("dashboard_bottom.png", Inches(2.5), Inches(4.8), height=Inches(2.2))
    except Exception as e:
        print("Warning: Images not found or error adding images:", e)


    # ----- Slide 6: Research and References -----
    slide6 = prs.slides.add_slide(blank_layout)
    txBox = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "HONEYWELL HACKATHON 2026\nRESEARCH AND REFERENCES"
    p.font.bold = True
    p.font.size = Pt(24)

    contentBox = slide6.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    tf2 = contentBox.text_frame
    tf2.word_wrap = True
    points = [
        "• Details / Links of the reference and research work:",
        "  - Liu, F. T., et al. (2008). \"Isolation Forest\". IEEE International Conference on Data Mining.",
        "  - Breiman, L. (2001). \"Random Forests\". Machine Learning, 45(1).",
        "  - MITRE ATT&CK Framework: Utilized for taxonomy mapping of simulated attack patterns.",
        "  - NIST SP 800-53: Guidelines on Security and Privacy Controls.",
        "",
        "• Technology Stack Documentation:",
        "  - scikit-learn (scikit-learn.org) - Core machine learning estimators.",
        "  - Plotly Dash (dash.plotly.com) - Real-time dashboard visualization.",
        "  - Pandas / NumPy (pandas.pydata.org) - Vectorized data pipeline transformations.",
        "",
        "• Open Source Assets:",
        "  - Developed during Honeywell Hackathon 2026 (MIT Licensed codebase)."
    ]
    for i, point in enumerate(points):
        p = tf2.add_paragraph()
        p.text = point
        if point.startswith("•"):
            p.font.bold = True
        p.font.size = Pt(16)

    # ----- Slide 7: Thank You -----
    slide7 = prs.slides.add_slide(blank_layout)
    txBox = slide7.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.bold = True
    p.font.size = Pt(48)
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "\nAI-Powered Behavioral Anomaly Detection for Cybersecurity\nHoneywell Hackathon 2026"
    p2.font.size = Pt(24)

    # Save presentation
    prs.save("Honeywell_Hackathon_Presentation.pptx")
    print("Successfully generated Honeywell_Hackathon_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
